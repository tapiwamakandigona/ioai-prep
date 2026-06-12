# Builds the Google Slides-ready PPTX from the rendered deck PDF.
# Each slide is a full-bleed image, so the design survives any import.
# Run:  uv run --with pymupdf --with python-pptx python build_pptx.py
import os, tempfile
import fitz
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, "..", "presentation", "how_machines_learn_to_see.pdf")
OUT = os.path.join(HERE, "..", "presentation", "how_machines_learn_to_see.pptx")

def main():
    doc = fitz.open(DECK)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as td:
        for i, page in enumerate(doc):
            png = os.path.join(td, f"{i}.png")
            page.get_pixmap(dpi=96).save(png)  # 1920x1080
            s = prs.slides.add_slide(blank)
            s.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(OUT)
    print("wrote", OUT, f"({len(doc)} slides)")

if __name__ == "__main__":
    main()
