#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # noqa: F401 – available for slide text alignment

INK = RGBColor(0x00, 0x00, 0x00)
DIM = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

S = [
 ("Vision Models: From Pixels to Diffusion",
  ["A tour of modern computer vision, from convolutions to image generation",
   "Built around the IOAI 2026 syllabus, Section 3: Computer Vision",
   "IOAI / GAITE preparation"],
  "Hi everyone. Today I'll walk you through how computers learn to see. We'll start with the basic question — what an image actually is to a computer — then build up through convolutional networks, detection and segmentation, transformers, and finally models that can generate brand-new images from noise. Everything here comes straight from the official IOAI 2026 syllabus for computer vision, so this is exactly the material I'm preparing for the olympiad."),
 ("An image is just a tensor of numbers",
  ["A color image is a grid of pixels: height × width × 3 channels (R, G, B), each value 0–255",
   "A small 224×224 photo is already ~150,000 numbers — meaning lives in patterns, not pixels",
   "The same object shifts, scales, rotates, changes lighting — the model must learn invariance",
   "Goal: compress raw pixels into a short, meaningful vector — an embedding"],
  "To a computer, a photo of a cat is not a cat — it's a giant grid of numbers, one brightness value per pixel per color channel. The hard part is that the meaning is spread across patterns of thousands of pixels, and those patterns move around: the cat can be anywhere in the frame, in any lighting. So vision models have one core job: compress this huge tensor of raw numbers into a short vector — an embedding — that captures what's in the image rather than where each pixel is."),
 ("Convolution: a small filter slides over the image  [THEORY]",
  ["A filter (kernel) is a tiny grid of learned weights, e.g. 3×3, slid across the whole image",
   "At each position: multiply filter weights by the pixels underneath, sum → one output number",
   "The output grid is a feature map — it lights up wherever the filter's pattern appears",
   "Weight sharing: same filter everywhere → few parameters; pattern found at any location",
   "Stacking layers detects edges → textures → parts → objects"],
  "This is the core idea of the whole field. Instead of connecting every pixel to every neuron, we use a tiny window of weights — say 3 by 3 — and slide it across the image. At each spot we multiply and sum, producing one number that says 'how strongly does this patch match my pattern?' The result is a feature map. Because the same filter is reused everywhere, we need very few parameters, and a pattern is found no matter where it sits. Early layers learn edges; deeper layers combine them into textures and object parts."),
 ("Pooling + stacking = a CNN for image classification",
  ["Max pooling: keep the largest value in each 2×2 window — halves resolution, keeps the strongest signal",
   "Average pooling: take the mean — used globally at the very end of most CNNs",
   "Pooling adds robustness to small shifts and shrinks the data flowing through the network",
   "Recipe: [conv → ReLU → pool] × N, then a linear layer → class probabilities (softmax)",
   "In code: nn.Conv2d, nn.MaxPool2d, nn.Linear — that's a full classifier"],
  "After convolution, we usually pool. Max pooling looks at each small window and keeps only the biggest value — 'was the feature here at all?' — which makes the network tolerant of the pattern shifting a pixel or two, and shrinks the feature maps so deeper layers are cheaper. Stack a few conv-and-pool blocks, flatten or average at the end, add one linear layer, and you have a complete image classifier: image in, probability per class out. In code, this is a few lines of PyTorch with nn.Conv2d, nn.MaxPool2d, and nn.Linear."),
 ("Augmentation: free training data from the data you have",
  ["Apply label-preserving changes during training: flip, random crop, noise, rotation, color jitter",
   "The model sees a 'new' image every epoch → less overfitting, better generalization",
   "Patching / random erasing: hide patches so the model can't rely on one local cue",
   "One line per transform (torchvision.transforms) — often the best accuracy boost in competitions",
   "Rule: the label must stay true — don't flip a '6' into a '9'"],
  "Labeled images are expensive, so we stretch what we have. During training we randomly flip, crop, add noise, or jitter the colors of each image — a flipped cat is still a cat, so the label survives, but the pixels are different. The model can't memorize exact pictures anymore and is forced to learn the real concept. Cutting out random patches goes further: the model must recognize a dog even when its head is hidden. In practice this is one line per transform in torchvision, and it's frequently the difference between an overfit model and a winning one."),
 ("ResNet: skip connections made deep networks work",
  ["Plain deep CNNs became hard to train — the training signal degrades through many layers",
   "ResNet's fix: each block outputs x + F(x) — the input plus a learned correction",
   "Gradients flow straight through the shortcut → 50–152-layer networks train reliably",
   "Trained on ImageNet, ResNet became the standard pre-trained encoder",
   "One line: torchvision.models.resnet50(weights=…) — reusable features for any task"],
  "For years, making CNNs deeper eventually made them worse — not from overfitting, but because the training signal degraded through so many layers. ResNet solved this with an almost embarrassingly simple idea: let each block learn only a correction to its input, and add the input back via a shortcut. If a block has nothing useful to add, it can pass the input through untouched. This made 50-plus-layer networks trainable, and ResNets trained on millions of ImageNet photos became the standard reusable 'vision backbone' everyone downloads instead of training from scratch."),
 ("Finetuning: start smart, not from scratch",
  ["Transfer learning: take a pre-trained encoder, swap the last layer for your classes",
   "Tiny data → freeze the backbone, train only the new head",
   "More data → finetune everything gently, with a low learning rate",
   "Parameter-efficient finetuning: train small inserted adapters (LoRA-style); backbone stays frozen",
   "Why it wins competitions: small datasets + pre-trained features beat training from zero"],
  "Here's the most practical slide of the talk. In a competition you might get two thousand labeled images — nowhere near enough to train a deep network from zero. So you don't. You download a model already trained on millions of images, chop off its last layer, and attach a new one for your classes. If data is tiny, freeze everything and train just that head. With a bit more data, unfreeze and finetune everything with a small learning rate. Parameter-efficient methods go further, training only small inserted adapter weights — cheaper and less likely to wreck the good pre-trained features."),
 ("Object detection: what is where — YOLO, SSD, DETR",
  ["Detection = classification + localization: a box, label, and confidence per object",
   "YOLO: one single pass predicts boxes over a grid — built for real-time speed",
   "SSD: single pass, predicts from feature maps at several scales → small and large objects",
   "DETR: a transformer outputs a direct set of detections — no anchor boxes, no NMS cleanup",
   "Pick: YOLO for speed · SSD classic multi-scale · DETR clean end-to-end"],
  "Classification says 'there's a dog somewhere.' Detection says 'a dog here, a ball there,' each with a rectangle and a confidence score. YOLO's idea is to do this in one single network pass over a grid of the image, which is why it runs in real time. SSD is also single-pass but predicts from several feature-map resolutions, so it catches both small and large objects. DETR replaces all the hand-engineered parts — anchor boxes, duplicate-removal — with a transformer that directly outputs a set of detections. In practice you'd load any of these pre-trained and finetune on your boxes."),
 ("Segmentation: a label for every pixel — U-Net",
  ["Output is a full-resolution mask: each pixel gets a class",
   "U-Net = encoder–decoder: encoder downsamples (what), decoder upsamples (where)",
   "Skip connections copy fine detail across the U → sharp mask boundaries",
   "Born in medical imaging; works well even with few training images",
   "Same skip idea as ResNet — here the skips carry spatial detail"],
  "Sometimes a box isn't enough — a medical model needs the exact outline of a tumor, pixel by pixel. That's segmentation. U-Net does it with a U-shaped network: the left side is a normal CNN encoder that shrinks the image while figuring out what's in it; the right side mirrors it, upsampling back to full resolution to decide where everything is. The trick is the skip connections across the U: they hand the decoder the fine details the encoder saw before downsampling, so the predicted mask has crisp, accurate edges instead of blurry blobs."),
 ("ViT: cut the image into patches, treat them like words  [THEORY]",
  ["Attention: each element scores its relevance to every other element, then gathers info by those weights",
   "Mechanic: each patch emits a query, compares it to all keys, mixes the values",
   "ViT: slice the image into 16×16 patches → embed each as a token → add position → transformer",
   "Global from layer one: patch 1 can directly attend to patch 196 — unlike a conv's local window",
   "Trade-off: no built-in locality bias → ViTs need large-scale pre-training to shine"],
  "Transformers conquered language first, and the question was: how do we feed an image to a model designed for word sequences? The Vision Transformer's answer: chop the image into a grid of small patches, flatten each into a vector — a token — and add a position signal so the model knows the layout. Then attention takes over: every patch asks, 'which other patches matter for understanding me?' and pulls in information accordingly, across the whole image at once. The catch is that ViTs don't have convolution's head start of assuming nearby pixels matter, so they need lots of pre-training data to win."),
 ("Self-supervised learning: features without labels",
  ["Human labels are expensive and slow; unlabeled images are nearly infinite",
   "Self-supervision: invent a training task from the data itself",
   "Contrastive idea: two augmented views of one image → pull together; other images → push apart",
   "Augmentations (slide 5) define what the model should ignore",
   "Result: a strong encoder from raw images — then finetune on your small labeled set"],
  "Everything so far assumed someone labeled the data — and labeling a million images is brutally expensive. Self-supervised learning sidesteps that: the data supervises itself. The most intuitive version is contrastive learning. Take one photo, make two random augmented versions — different crops, flips, colors. Tell the network: these two must land close together in embedding space, while views of other images get pushed away. To succeed, the network must learn what makes that image itself, ignoring superficial changes. You get a powerful encoder from unlabeled data, then finetune it with the few labels you have."),
 ("CLIP: images and text in one shared space",
  ["Two encoders — image and text — trained on huge sets of (image, caption) pairs",
   "Contrastive: matching pairs → similar embeddings; mismatched → dissimilar",
   "Result: 'a photo of a dog' (text) lands next to actual dog photos",
   "Zero-shot classification: embed candidate captions + the image, pick the closest — no training",
   "Also powers image search and guides text-to-image generators"],
  "CLIP takes the contrastive idea from the last slide and applies it across two worlds at once. One encoder turns images into vectors; another turns text into vectors; training pulls each image toward its true caption and away from all wrong ones, over hundreds of millions of pairs. The payoff is a shared space where words and pictures are directly comparable. That enables something remarkable: zero-shot classification. Want a cat-versus-dog classifier with zero training images? Embed the sentences 'a photo of a cat' and 'a photo of a dog,' embed your image, and pick the nearest sentence. It just works, and it's a go-to tool in olympiad tasks."),
 ("GANs: a forger and a detective, trained against each other",
  ["Generator turns random noise into images; discriminator judges real vs. fake",
   "Adversarial loop: better detection forces better forgery — both improve together",
   "At equilibrium, fakes become hard to distinguish from real photos",
   "Strengths: sharp images, one-pass fast generation",
   "Pain points: unstable training, mode collapse (only a few 'safe' outputs)"],
  "Now we flip the direction: instead of understanding images, we create them. A GAN is a game between two networks. The generator is a forger: it takes random noise and tries to paint a convincing image. The discriminator is a detective: shown real photos and forgeries, it learns to tell them apart. Each one's progress forces the other to improve — better detection demands better forgery. When training goes well, the fakes become indistinguishable. The catch is that this two-player game is unstable: it can oscillate, or the generator can collapse to producing the same few images over and over."),
 ("Diffusion: learn to remove noise, then start from pure noise",
  ["Forward (fixed): add a little noise to a real image, step by step, until pure static",
   "Reverse (learned): a network predicts the noise so each step can subtract it",
   "Generation: start from random static, denoise step by step → a new image condenses out",
   "Condition on text embeddings (CLIP-style) → text-to-image, e.g. Stable Diffusion",
   "vs. GANs: stabler training, more diverse outputs — but many steps, so slower"],
  "Diffusion models are today's state of the art for image generation, and the idea is beautifully simple. First, destroy: take real images and add a bit of noise, again and again, until nothing but static remains — that part needs no learning. Then teach a network to undo one small step of that noising: given a noisy image, predict the noise so we can subtract it. To generate, hand the model pure random static and let it denoise step by step — an image gradually condenses out of the noise. Condition each step on a text embedding and you get text-to-image generation, like Stable Diffusion."),
 ("Where each topic sits in the IOAI 2026 syllabus",
  ["Convolutional layers — THEORY + PRACTICE (slide 3)",
   "Image classification + pooling (4) · Augmentation (5) · Pre-trained encoders / ResNet (6)",
   "Model finetuning (7) · Detection: YOLO, SSD, DETR (8) · Segmentation: U-Net (9)",
   "Attention / Transformers → ViT — supporting THEORY (10) · Self-supervised learning (11)",
   "CLIP (12) · GANs (13) · Diffusion models (14)"],
  "Quick orientation: this is the official checklist. In IOAI terms, 'Theory' means I must be able to explain how something works on paper — that's convolutions, and the attention mechanism behind vision transformers. 'Practice' means I need working intuition: what the tool does, when to reach for it, and how to use it with standard libraries like PyTorch, torchvision, and Hugging Face — typically loading something pre-trained and finetuning it rather than building it from scratch. Every row on this table maps to one slide you've just seen."),
 ("From pixels to diffusion — and what I'd build next",
  ["One storyline: pixels → conv → ResNet → detect/segment → ViT → self-supervised → CLIP → GAN/diffusion",
   "Two theory anchors to master cold: convolution and attention",
   "One winning strategy: pre-trained encoder + strong augmentation + careful finetuning",
   "Next project: finetune ResNet on a small custom dataset; compare against CLIP zero-shot",
   "Questions welcome — the study guide has the details"],
  "So, the whole field in one breath: images are tensors; convolutions extract local patterns cheaply; ResNets made depth trainable and gave us reusable encoders; detection and segmentation localize what classifiers can only name; transformers brought global attention to vision via patches; self-supervision and CLIP removed the labeling bottleneck; and GANs and diffusion turned the pipeline around to generate images. My next step is hands-on: finetune a pre-trained ResNet on a small dataset with strong augmentation, and compare it against CLIP's zero-shot predictions as a baseline. Thanks — happy to take questions."),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for idx, (title, bullets, notes) in enumerate(S):
    slide = prs.slides.add_slide(blank)
    is_cover = idx == 0
    # eyebrow
    eb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.4))
    p = eb.text_frame.paragraphs[0]
    p.text = "IOAI 2026 · COMPUTER VISION" if is_cover else f"VISION MODELS · {idx+1:02d} / 16"
    p.font.size = Pt(11); p.font.color.rgb = DIM; p.font.name = "Consolas"
    # title
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.4) if is_cover else Inches(0.95),
                                  Inches(12), Inches(2.6) if is_cover else Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(54) if is_cover else Pt(32)
    p.font.bold = True; p.font.color.rgb = INK
    # bullets
    by = Inches(4.6) if is_cover else Inches(2.35)
    bb = slide.shapes.add_textbox(Inches(0.7), by, Inches(11.9), Inches(4.6))
    tf = bb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("" if is_cover else "—  ") + b
        p.font.size = Pt(16) if is_cover else Pt(18)
        p.font.color.rgb = DIM if is_cover else INK
        p.space_after = Pt(14)
    # speaker notes
    slide.notes_slide.notes_text_frame.text = notes

prs.save("/work/projects/ioai-prep/vision-models/vision_models_deck.pptx")
print("pptx saved,", len(S), "slides")
